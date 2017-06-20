#!/usr/bin/env python
# encoding: utf-8

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION
import time
t = time.time()
prs = Presentation()

for i in range(1, 30):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 在幻灯片中加入一页6号风格（空白）幻灯片

    # chart1 左上方图
    x, y, cx, cy = Inches(0.5), Inches(0.5), Inches(4), Inches(3)  # 按英尺标准指定x，y值

    chart_data = ChartData()  # 图表data类

    chart_data.categories = [u'A班级得分率', u'B班级得分率']  # 图表加入两栏
    chart_data.add_series(u'得分率对比', (80.5, 60.5))  # 在两栏分别填入数据

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )  # add_chart(图表类型，xy表示图表位置，cx cy表示图表宽高，并且插入chart_data中规定好的数据）

    chart = graphic_frame.chart  # 从生成的图表中取出图表类
    chart.chart_style = 21  # 图表整体颜色风格

    chart.has_title = True  # 图表是否含有标题，默认为False
    chart.chart_title.text_frame.clear()  # 清除原标题
    new_paragraph = chart.chart_title.text_frame.add_paragraph()  # 添加一行新标题
    new_paragraph.text = '得分率对比'  # 新标题
    new_paragraph.font.size = Pt(15)  # 新标题字体大小

    category_axis = chart.category_axis  # category_axis 为chart的category控制类
    category_axis.has_major_gridlines = True  # 是否显示纵轴线
    category_axis.tick_labels.font.italic = True  # tick_labels为图表下标签，置为斜体
    category_axis.tick_labels.font.size = Pt(15)  # 下标签字体大小
    category_axis.tick_labels.font.color.rgb = RGBColor(255, 0, 0)  # 标签字体颜色

    value_axis = chart.value_axis  # value_axis 为chart的value控制类
    value_axis.maximum_scale = 100.0  # 纵坐标最大值
    value_axis.minimum_scale = 0.0  # 纵坐标最小值
    value_axis.minor_tick_mark = XL_TICK_MARK.CROSS
    value_axis.has_minor_gridlines = True

    tick_labels = value_axis.tick_labels  # tick_labels 为chart的纵轴标签控制类
    tick_labels.number_format = '0"%"'  # 标签显示样式
    tick_labels.font.bold = True  # 字体加粗
    tick_labels.font.size = Pt(14)  # 字体大小
    tick_labels.font.color.rgb = RGBColor(0, 255, 0)  # 标签颜色

    plot = chart.plots[0]  # 取图表中第一个plot
    plot.has_data_labels = True  # 是否显示数据标签
    data_labels = plot.data_labels  # 数据标签控制类
    data_labels.font.size = Pt(13)  # 字体大小
    data_labels.font.color.rgb = RGBColor(0, 0, 255)  # 字体颜色
    data_labels.position = XL_LABEL_POSITION.INSIDE_END  # 字体位置

    # chart 2 左下方图
    x, y, cx, cy = Inches(0.5), Inches(3.5), Inches(4), Inches(3)  # 按英尺标准指定x，y值
    chart_data = ChartData()
    chart_data.categories = ['A', 'B', 'C', 'D']
    chart_data.add_series(u'A班级选项占比', (80, 10, 9, 10))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart  # PIE为饼状图

    chart.has_legend = True  # 是否含有下方的说明
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.horz_offset = 0  # 说明位移量 [-1, 1] 默认为0

    chart.plots[0].has_data_labels = True  # 饼中是否写入数值
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'  # 数值显示格式
    data_labels.position = XL_LABEL_POSITION.INSIDE_END  # 数值布局方式

    chart.has_title = True
    chart.chart_title.text_frame.clear()  # 清除原标题
    new_paragraph = chart.chart_title.text_frame.add_paragraph()  # 添加一行新标题
    new_paragraph.text = 'A班级选项占比'  # 新标题
    new_paragraph.font.size = Pt(13)  # 新标题字体大小

    # chart 3 右下方图
    x, y, cx, cy = Inches(5.5), Inches(4), Inches(4), Inches(3)  # 按英尺标准指定x，y值
    chart_data = ChartData()
    chart_data.categories = ['A', 'B', 'C', 'D']
    chart_data.add_series(u'B班级选项占比', (0.1, 0.2, 0.3, 0.4))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_LABEL_POSITION.INSIDE_END

    chart.has_title = True
    chart.chart_title.text_frame.clear()  # 清除原标题
    new_paragraph = chart.chart_title.text_frame.add_paragraph()  # 添加一行新标题
    new_paragraph.text = 'B班级选项占比'  # 新标题
    new_paragraph.font.size = Pt(13)  # 新标题字体大小


    # chart 4 右上方图
    x, y, cx, cy = Inches(5.5), Inches(0.5), Inches(4), Inches(3)
    chart_data = ChartData()
    chart_data.categories = ['0', '1-3', '4-6', '7-9']
    chart_data.add_series('', (50, 18, 30, 34))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(13)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_LABEL_POSITION.INSIDE_END

    chart.has_title = True
    chart.chart_title.text_frame.clear()
    new_title = chart.chart_title.text_frame.add_paragraph()
    new_title.text = '得分占比'
    new_title.font.size = Pt(13)

    slide.shapes.add_picture('img_path.jpg', Inches(1), Inches(1))
prs.save('test.pptx')
print time.time() - t